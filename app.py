from flask import Flask, request, jsonify
from dotenv import load_dotenv
import uuid
from PyPDF2 import PdfReader
from pptx import Presentation
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tempfile
import google.generativeai as genai
from langchain_experimental.agents import create_csv_agent
from langchain.llms import OpenAI
from langchain.document_loaders import PDFPlumberLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Chroma
from langchain.chains import RetrievalQA
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
import pandas as pd
import io
import tempfile
from PIL import Image
import pymongo
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_experimental.agents.agent_toolkits import create_csv_agent
from langchain.prompts import PromptTemplate
import requests
from bs4 import BeautifulSoup
import speech_recognition as sr
from googletrans import Translator,LANGUAGES
from werkzeug.utils import secure_filename
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)


load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
model = genai.GenerativeModel("gemini-1.5-flash-latest")

app = Flask(__name__)

def load_excel_and_convert_to_csv(file, query):
    if file:
        with tempfile.NamedTemporaryFile(mode="w+", suffix=".csv", delete=False) as f:
            data_str = file.read().decode("utf-8")
            f.write(data_str)
            f.flush()
            model = ChatGoogleGenerativeAI(model="gemini-1.5-flash-latest", temperature=0.8)
            agent = create_csv_agent(model, f.name, verbose=True)
            response = agent.run(query)
            return response

def get_ppt_content(ppt_docs):
    slides_content = []
    for ppt in ppt_docs:
        prs = Presentation(ppt)
        for slide in prs.slides:
            slide_text = ""
            slide_images = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as img_file:
                        img_file.write(image_bytes)
                        img_file.flush()
                        slide_images.append(img_file.name)
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            slides_content.append((slide_text, slide_images))
    return slides_content

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks, vector_store_path):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local(vector_store_path)

def load_vector_store(vector_store_path):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    return FAISS.load_local(vector_store_path, embeddings, allow_dangerous_deserialization=True)

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details,and also sometimes the question might not always be in english so u gotta convert it to english and then find the necessay answer provided to you in the context, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-1.5-flash-latest", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question, vector_store_path):
    new_db = load_vector_store(vector_store_path)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

def get_gemini_response1(input_text, image):
    model = genai.GenerativeModel('gemini-pro-vision')
    if input_text:
        response = model.generate_content([input_text, image])
    else:
        response = model.generate_content(image)
    return response.text

def fetch_related_links(query):
    search_url = "https://www.googleapis.com/customsearch/v1"
    params = {
        'key': os.getenv('GOOGLE_KEY'),
        'cx': os.getenv('SEARCH_ENGINE_ID'),
        'q': query
    }
    response = requests.get(search_url, params=params)
    search_results = response.json()
    links = [item['link'] for item in search_results.get('items', [])]
    return links
def fetch_images(query):
    search_url = "https://www.google.com/search"
    params = {
        'q': query,
        'tbm': 'isch'
    }
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3'}
    response = requests.get(search_url, params=params, headers=headers)
    soup = BeautifulSoup(response.text, 'html.parser')
    images = [img['src'] for img in soup.find_all('img') if img.has_attr('src')]
    # Filter out non-image URLs
    images = [img for img in images if img.startswith('http')]
    return images

@app.route('/api/excel', methods=['POST'])
def excel_qa():
    file = request.files.get('file')
    query = request.form.get('query')
    response = load_excel_and_convert_to_csv(file, query)
    return jsonify({'response': response})

@app.route('/api/pdf', methods=['POST'])
def pdf_qa():
    files = request.files.getlist('file')
    user_question = request.form.get('question')
    raw_text = get_pdf_text(files)
    
    if raw_text:
        text_chunks = get_text_chunks(raw_text)
        get_vector_store(text_chunks, "faiss_index_pdf")
        response = user_input(user_question, "faiss_index_pdf")
        return jsonify({'response': response})
    else:
        return jsonify({"error": "No text extracted from PDF files."}), 400

@app.route('/api/ppt', methods=['POST'])
def ppt_qa():
    user_question = request.form.get('question')
    ppt_docs = request.files.getlist('file')
    translator = Translator()
    try:
        print(user_question)
        detected_lang = translator.detect(user_question).lang
        print(f"Detected language: {LANGUAGES.get(detected_lang, 'Unknown')}")
        if detected_lang != 'en':
            translated_question = translator.translate(user_question, src=detected_lang, dest='en').text
            print(f"Translated question: {translated_question}")
        else:
            translated_question = user_question
    except Exception as e:
        print(f"Error in translation: {e}")
        return jsonify({"error": "Translation service failed."}), 500
    if ppt_docs:
        slides_content = get_ppt_content(ppt_docs)
        raw_text = "\n".join([content[0] for content in slides_content])
        text_chunks = get_text_chunks(raw_text)
        get_vector_store(text_chunks, "faiss_index_ppt")
        response_text = user_input(translated_question, "faiss_index_ppt")
        slide_images = []
        for content in slides_content:
            if user_question.lower() in content[0].lower():
                slide_images.extend(content[1])
        related_links = fetch_related_links(response_text.split('\n')[0])
        relevant_images = fetch_images(response_text.split('\n')[0])

        
        response = {
            "text": response_text,
            "images_from_slides": slide_images,
            "images_from_internet": relevant_images,
            "links": related_links
        }
        return jsonify(response)
    else:
        return jsonify({"error": "No PPT files uploaded."}), 400
@app.route('/ppt/audio', methods=['POST'])
def load_audio():
    file = request.files['audio']
    ppt_docs = request.files.getlist('file')
    filename = secure_filename(file.filename)
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(file_path)

    recognizer = sr.Recognizer()
    with sr.AudioFile(file_path) as source:
        audio_data = recognizer.record(source)
        transcription = recognizer.recognize_google(audio_data, show_all=True)

    detected_language = transcription['alternative'][0]['language']
    transcribed_text = transcription['alternative'][0]['transcript']

    if detected_language != 'en':
        translator = Translator()
        translated = translator.translate(transcribed_text, src=detected_language, dest='en')
        transcribed_text = translated.text
    if ppt_docs:
        slides_content = get_ppt_content(ppt_docs)
        raw_text = "\n".join([content[0] for content in slides_content])
        text_chunks = get_text_chunks(raw_text)
        get_vector_store(text_chunks, "faiss_index_ppt")
        response_text = user_input(transcribed_text, "faiss_index_ppt")
        slide_images = []
        for content in slides_content:
            if transcribed_text.lower() in content[0].lower():
                slide_images.extend(content[1])
        related_links = fetch_related_links(response_text.split('\n')[0])
        relevant_images = fetch_images(response_text.split('\n')[0])

        
        response = {
            "text": response_text,
            "images_from_slides": slide_images,
            "images_from_internet": relevant_images,
            "links": related_links
        }
        return jsonify(response)
    else:
        return jsonify({"error": "No PPT files uploaded."}), 400

@app.route('/api/vision', methods=['POST'])
def vision():
    input_text = request.form.get('input')
    image_file = request.files.get('image')
    image = Image.open(image_file)
    response = get_gemini_response1(input_text, image)
    return jsonify({'response': response})

if __name__ == '__main__':
    app.run(port=5000, debug=True)